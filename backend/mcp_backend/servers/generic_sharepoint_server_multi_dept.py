#!/usr/bin/env python3
"""
Multi-Department Generic SharePoint MCP Server with On-Behalf-Of (OBO) Flow Support

This server provides SharePoint document search and retrieval capabilities
for multiple departments/organizations using Microsoft Graph API with OBO flow for user-scoped access.

This is a generic, reusable implementation that can be instantiated for any department
by providing a department prefix (e.g., MPO, FIN, IT) to access department-specific environment variables.

Global environment variables (shared across all departments):
- SHP_USE_DELEGATED_ACCESS: Enable user-scoped access (true/false, default: true)
- SHP_OBO_SCOPE: Scopes for On-Behalf-Of token exchange (default: standard Graph scopes)
- SHP_GLOBAL_SEARCH_LIMIT: Maximum total results across all folders before terminating search (default: 30)
- SHP_FOLDER_SEARCH_LIMIT: Maximum number of discovered folders to search (default: 5)

Department-specific environment variables (prefix with DEPT_):
- {DEPT}_SHP_ID_APP: Azure App Registration Client ID
- {DEPT}_SHP_ID_APP_SECRET: Azure App Registration Client Secret
- {DEPT}_SHP_TENANT_ID: Azure AD Tenant ID
- {DEPT}_SHP_SITE_URL: SharePoint Site URL
- {DEPT}_SHP_ORG_NAME: Organization name for tool descriptions
- {DEPT}_SHP_DOC_LIBRARY: SharePoint Document Library path (optional)
- {DEPT}_SHP_DEFAULT_SEARCH_FOLDERS: Comma-separated list of default folders to search in

Example usage:
- For MPO: create server with department_prefix="MPO"
- For Finance: create server with department_prefix="FIN"
- For IT: create server with department_prefix="IT"
"""

import os
import sys
import asyncio
import json
import logging
import re
import zipfile
import xml.etree.ElementTree as ET
from typing import Dict, List, Any, Optional
from pathlib import Path
from io import BytesIO
from concurrent.futures import ThreadPoolExecutor, as_completed
import threading
from functools import lru_cache

# Add local modules to path
sys.path.insert(0, str(Path(__file__).parent))

try:
    from sharepoint_oauth_client import SharePointOAuthClient
    from fastmcp import FastMCP

    mcp = FastMCP("Generic SharePoint Server - Multi-Department")

    # Global MCP server instance (will be configured per department)
    generic_sharepoint_server = mcp

except ImportError as e:
    logging.error(f"Failed to import required modules: {e}")
    raise


# Authentication error message for delegated access mode
DELEGATED_ACCESS_AUTH_ERROR = (
    "DELEGATED ACCESS MODE: User authentication token is required but not provided. "
    "SharePoint access with user permissions requires a valid OAuth2 access token. "
    "This request cannot proceed without proper authentication. "
    "Either provide a valid user token or set SHP_USE_DELEGATED_ACCESS=false to use application access."
)


def clean_env_var(value: str) -> str:
    """Clean environment variable by removing quotes and inline comments"""
    if not value:
        return value
    # Remove quotes and inline comments
    value = value.strip("\"'")
    value = re.sub(r"\s*#.*$", "", value)
    value = value.strip("\"'").strip()
    return value


class DepartmentSharePointConfig:
    """Configuration class for department-specific SharePoint settings"""

    def __init__(self, department_prefix: str):
        self.department_prefix = department_prefix.upper()

        # Global settings (shared across all departments)
        # Delegated access is always enabled - no application fallback
        self.use_delegated_access = True
        self.obo_scope = clean_env_var(
            os.getenv(
                "SHP_OBO_SCOPE",
                "https://graph.microsoft.com/Sites.Read.All https://graph.microsoft.com/Files.Read.All",
            )
        )

        # Performance settings - configurable via environment variables
        self.global_search_limit = int(
            clean_env_var(os.getenv("SHP_GLOBAL_SEARCH_LIMIT", "30"))
        )
        self.folder_search_limit = int(
            clean_env_var(os.getenv("SHP_FOLDER_SEARCH_LIMIT", "5"))
        )

        # Department-specific settings
        self.site_url = clean_env_var(
            os.getenv(f"{self.department_prefix}_SHP_SITE_URL", "")
        )
        self.org_name = clean_env_var(
            os.getenv(
                f"{self.department_prefix}_SHP_ORG_NAME",
                f"{department_prefix} Department",
            )
        )
        self.doc_library = clean_env_var(
            os.getenv(f"{self.department_prefix}_SHP_DOC_LIBRARY", "")
        )
        self.default_search_folders = (
            os.getenv(f"{self.department_prefix}_SHP_DEFAULT_SEARCH_FOLDERS", "").split(
                ","
            )
            if os.getenv(f"{self.department_prefix}_SHP_DEFAULT_SEARCH_FOLDERS")
            else []
        )

        # Clean up empty strings and comments in search folders
        self.default_search_folders = [
            clean_env_var(folder)
            for folder in self.default_search_folders
            if clean_env_var(folder)
        ]

        # Validation
        self.validate_config()

    def validate_config(self):
        """Validate that required configuration is present"""
        # Only site URL is required for delegated access mode
        if not self.site_url:
            raise ValueError(
                f"Missing required environment variable: {self.department_prefix}_SHP_SITE_URL"
            )


# Global variables (will be set when server is initialized)
config: Optional[DepartmentSharePointConfig] = None
oauth_client: Optional[SharePointOAuthClient] = None
logger = logging.getLogger(__name__)


def initialize_department_server(department_prefix: str):
    """Initialize the server for a specific department"""
    global config, oauth_client

    try:
        # Load department-specific configuration
        config = DepartmentSharePointConfig(department_prefix)

        # Initialize OAuth client with delegated access (credentials are placeholders)
        oauth_client = SharePointOAuthClient(
            site_url=config.site_url,
            use_delegated_access=config.use_delegated_access,
            obo_scope=config.obo_scope,
        )

        logger.info(f"Initialized SharePoint server for department: {config.org_name}")
        logger.info(
            f"Configuration: delegated_access={config.use_delegated_access}, site_url={config.site_url}"
        )

        return True

    except Exception as e:
        logger.error(
            f"Failed to initialize SharePoint server for {department_prefix}: {e}"
        )
        # Return False to allow other MCP servers to continue
        return False


def extract_user_token(context: Optional[Dict[str, Any]] = None) -> Optional[str]:
    """Extract the user's OAuth token from context or environment"""
    # First try context (passed from crew MCP manager)
    if context and "user_token" in context:
        return context["user_token"]

    # Fallback to environment variable
    user_token = os.getenv("USER_JWT_TOKEN")
    if user_token and user_token != "user_token_placeholder":
        return user_token

    return None


async def discover_matching_folders(
    oauth_client,
    access_token: str,
    site_id: str,
    drive_id: str,
    folder_pattern: str,
    max_depth: int = 2,
) -> List[str]:
    """
    Dynamically discover folders that match a given pattern.

    This function searches through SharePoint folder structure to find folders
    that match the specified pattern (e.g., "Major Projects Office", "MPO").

    Args:
        oauth_client: SharePoint OAuth client instance
        access_token: OAuth access token
        site_id: SharePoint site ID
        drive_id: SharePoint drive ID
        folder_pattern: Pattern to match (can be partial folder names)
        max_depth: Maximum depth to search

    Returns:
        List of folder paths that match the pattern
    """
    logger.info(f"Discovering folders matching pattern: {folder_pattern}")

    # Parse the pattern to extract key terms dynamically
    pattern_parts = folder_pattern.lower().split("/")
    target_folder = pattern_parts[-1] if pattern_parts else folder_pattern.lower()

    # Extract key terms dynamically from the target folder pattern
    key_terms = []

    # Split on common separators and extract meaningful terms
    terms = re.split(r"[\s\-_()]+", target_folder)
    key_terms = [term.strip() for term in terms if len(term) > 2]

    # Also include common abbreviations if we can detect them
    # Look for potential abbreviations (uppercase sequences or known patterns)
    for term in terms:
        term = term.strip()
        if len(term) > 1:
            # Add the term as-is
            if term not in key_terms:
                key_terms.append(term)

            # If it's a multi-word phrase, add individual words
            if " " in term:
                individual_words = term.split()
                for word in individual_words:
                    if len(word) > 2 and word not in key_terms:
                        key_terms.append(word)

    logger.info(f"Extracted key terms for folder discovery: {key_terms}")

    matching_folders = []

    async def search_folders_recursive(current_path: str, depth: int = 0):
        if depth >= max_depth:
            return

        try:
            # Get items from current folder
            success, items_data = await oauth_client.get_drive_items(
                access_token, site_id, drive_id, current_path
            )

            if not success:
                logger.warning(f"Failed to access folder for discovery: {current_path}")
                return

            items = items_data.get("value", [])
            folder_items = [item for item in items if item.get("folder")]

            for folder_item in folder_items:
                folder_name = folder_item.get("name", "")
                folder_name_lower = folder_name.lower()

                # Check if this folder matches our pattern
                match_score = 0
                for term in key_terms:
                    if term in folder_name_lower:
                        match_score += 1

                # If we have a good match (at least half the key terms)
                if match_score >= len(key_terms) // 2 and match_score > 0:
                    folder_path = (
                        f"{current_path}/{folder_name}" if current_path else folder_name
                    )
                    matching_folders.append(folder_path)
                    logger.info(
                        f"Found matching folder: {folder_path} (score: {match_score})"
                    )

                # Continue searching in subfolders
                if depth < max_depth - 1:
                    subfolder_path = (
                        f"{current_path}/{folder_name}" if current_path else folder_name
                    )
                    await search_folders_recursive(subfolder_path, depth + 1)

        except Exception as e:
            logger.warning(f"Error during folder discovery at {current_path}: {e}")

    # Start discovery from root
    await search_folders_recursive("", 0)

    if matching_folders:
        logger.info(
            f"Discovered {len(matching_folders)} matching folders: {matching_folders}"
        )
    else:
        logger.warning(f"No folders found matching pattern: {folder_pattern}")
        # Fallback: try using the pattern as-is
        matching_folders = [folder_pattern]

    return matching_folders


async def search_folder_recursive(
    oauth_client,
    access_token: str,
    site_id: str,
    drive_id: str,
    folder_path: str,
    query: str,
    org_name: str,
    max_depth: int = 2,
    current_depth: int = 0,
    max_results: int = 20,
    max_folders_per_level: int = 10,
) -> Dict[str, Any]:
    """
    Recursively search a SharePoint folder and all its subfolders for documents matching the query.

    Args:
        oauth_client: SharePoint OAuth client instance
        access_token: OAuth access token
        site_id: SharePoint site ID
        drive_id: SharePoint drive ID
        folder_path: Folder path to search
        query: Search query
        org_name: Organization name for metadata
        max_depth: Maximum recursion depth to prevent infinite loops (default: 2)
        current_depth: Current recursion depth
        max_results: Maximum number of results to return before early termination (default: 20)
        max_folders_per_level: Maximum number of folders to search per depth level (default: 10)

    Returns:
        Dict containing matching documents, total files, and folders searched
    """
    logger.info(f"Searching folder recursively: {folder_path} (depth: {current_depth})")

    if current_depth >= max_depth:
        logger.warning(f"Maximum depth {max_depth} reached for folder: {folder_path}")
        return {
            "matching_docs": [],
            "total_files": 0,
            "folders_searched": [folder_path],
        }

    try:
        # Get items from current folder
        success, items_data = await oauth_client.get_drive_items(
            access_token, site_id, drive_id, folder_path
        )

        if not success:
            logger.warning(f"Failed to access folder: {folder_path} - {items_data}")
            return {
                "matching_docs": [],
                "total_files": 0,
                "folders_searched": [folder_path],
            }

        items = items_data.get("value", [])

        # Separate files and folders
        file_items = [item for item in items if item.get("file")]
        folder_items = [item for item in items if item.get("folder")]

        # Limit number of folders to search per level to prevent excessive recursion
        if len(folder_items) > max_folders_per_level:
            folder_items = folder_items[:max_folders_per_level]
            logger.info(
                f"Limited folder search to {max_folders_per_level} folders at depth {current_depth}"
            )

        matching_docs = []
        total_files = len(file_items)
        folders_searched = [folder_path]

        # Search files in current folder
        query_words = query.lower().split()

        for item in file_items:
            name = item.get("name", "")
            name_lower = name.lower()

            # Match if any significant word from query appears in filename
            match_found = False
            for word in query_words:
                if len(word) > 2 and word in name_lower:  # Skip very short words
                    match_found = True
                    break

            # Also check for exact phrase match
            if query.lower() in name_lower:
                match_found = True

            if match_found:
                matching_docs.append(
                    {
                        "name": name,
                        "id": item.get("id"),
                        "size": item.get("size"),
                        "lastModified": item.get("lastModifiedDateTime"),
                        "webUrl": item.get("webUrl"),
                        "folder_path": folder_path,
                        "organization": org_name,
                    }
                )

                # Early termination within file loop if we have enough results from this folder
                if len(matching_docs) >= max_results:
                    logger.info(
                        f"Early termination in file search: Found {len(matching_docs)} results in folder {folder_path}"
                    )
                    break

        # Early termination if we have enough results
        if len(matching_docs) >= max_results:
            logger.info(
                f"Early termination: Found {len(matching_docs)} results, stopping search"
            )
            return {
                "matching_docs": matching_docs,
                "total_files": total_files,
                "folders_searched": folders_searched,
            }

        # Recursively search subfolders
        for subfolder in folder_items:
            subfolder_name = subfolder.get("name", "")

            # Skip system folders and hidden folders
            if subfolder_name.startswith(".") or subfolder_name.startswith("~"):
                continue

            subfolder_path = (
                f"{folder_path}/{subfolder_name}" if folder_path else subfolder_name
            )

            # Recursive call
            subfolder_results = await search_folder_recursive(
                oauth_client,
                access_token,
                site_id,
                drive_id,
                subfolder_path,
                query,
                org_name,
                max_depth,
                current_depth + 1,
                max_results
                - len(matching_docs),  # Reduce max_results for deeper levels
                max_folders_per_level,
            )

            # Combine results
            matching_docs.extend(subfolder_results["matching_docs"])
            total_files += subfolder_results["total_files"]
            folders_searched.extend(subfolder_results["folders_searched"])

            # Early termination if we have enough results after processing this subfolder
            if len(matching_docs) >= max_results:
                logger.info(
                    f"Early termination in subfolder loop: Found {len(matching_docs)} total results, stopping subfolder search"
                )
                break

        return {
            "matching_docs": matching_docs,
            "total_files": total_files,
            "folders_searched": folders_searched,
        }

    except Exception as e:
        logger.error(f"Error in recursive search for folder {folder_path}: {e}")
        return {
            "matching_docs": [],
            "total_files": 0,
            "folders_searched": [folder_path],
        }


@generic_sharepoint_server.tool()
async def search_sharepoint_documents(
    query: str = "documents",
    folder_path: Optional[str] = None,
    search_subfolders: bool = True,
    user_token: Optional[str] = None,
) -> Dict[str, Any]:
    """
    Search for documents in SharePoint based on a query using Microsoft Graph API with OBO.

    This tool searches for documents in SharePoint folders with user-scoped access.
    It can search in a specific folder or across multiple configured default folders.

    Args:
        query: Search terms to find relevant documents (e.g., "strategy", "project")
        folder_path: Specific SharePoint folder path to search in (optional)
        search_subfolders: Whether to search in subfolders as well (default: True)
        user_token: Optional user access token for delegated access

    Returns:
        JSON response with search results and document list
    """
    if not config or not oauth_client:
        return {
            "status": "error",
            "error": "Server not initialized with department configuration",
            "message": "Please initialize server with a valid department prefix",
        }

    try:
        # Get user token from environment variable (set by crew MCP manager) if not provided
        if not user_token:
            user_token = os.getenv("USER_JWT_TOKEN")

        # Validate we have a proper user token for delegated access
        if (
            user_token == "user_token_placeholder"
            or not user_token
            or not user_token.strip()
        ):
            logger.error(
                "No valid OAuth token available. SharePoint access requires user authentication."
            )
            return {
                "status": "error",
                "error": "Authentication required",
                "message": "You must be authenticated to access SharePoint documents. Please ensure you're logged in and have the necessary permissions.",
                "organization": config.org_name if config else "Unknown",
            }

        logger.info(f"Using OAuth token for SharePoint OBO flow")
        # Log token type for debugging
        if user_token and len(user_token) > 20:
            logger.debug(f"Token type check - starts with: {user_token[:20]}...")

        logger.info(f"Searching {config.org_name} SharePoint for: {query}")

        # Get appropriate access token using OBO flow
        access_token = await oauth_client.get_access_token(user_token)
        if not access_token:
            return {
                "status": "error",
                "query": query,
                "error": "Failed to get access token",
                "message": f"Unable to authenticate with {config.org_name} SharePoint. This may be due to insufficient permissions or expired credentials. Please ensure you have the necessary SharePoint access permissions for this organization.",
                "organization": config.org_name,
            }

        # Get site and drive information
        success, site_data = await oauth_client.get_site_info(access_token)
        if not success:
            return {
                "status": "error",
                "query": query,
                "error": site_data,
                "message": "Failed to access SharePoint site",
                "organization": config.org_name,
            }

        site_id = site_data.get("id")

        # Get drives (document libraries)
        success, drives_data = await oauth_client.get_site_drives(access_token, site_id)
        if not success:
            return {
                "status": "error",
                "query": query,
                "error": drives_data,
                "message": "Failed to access document libraries",
                "organization": config.org_name,
            }

        drives = drives_data.get("value", [])
        if not drives:
            return {
                "status": "error",
                "query": query,
                "error": "No document libraries found",
                "message": "No accessible document libraries found",
                "organization": config.org_name,
            }

        # Use first drive for now
        drive_id = drives[0].get("id")

        # Dynamic folder discovery: find folders that match the configured patterns
        discovered_folders = []

        if folder_path:
            # If specific folder path provided, use it directly
            discovered_folders = [folder_path]
        elif config.default_search_folders:
            # Discover actual folders that match the configured patterns
            for pattern in config.default_search_folders:
                matching_folders = await discover_matching_folders(
                    oauth_client, access_token, site_id, drive_id, pattern
                )
                discovered_folders.extend(matching_folders)

        # If no folders discovered, try root level
        if not discovered_folders:
            discovered_folders = [""]

        # Limit the number of discovered folders to search to prevent excessive search time
        if len(discovered_folders) > config.folder_search_limit:
            discovered_folders = discovered_folders[: config.folder_search_limit]
            logger.info(
                f"Limited folder discovery to first {config.folder_search_limit} folders to improve performance"
            )

        all_matching_documents = []
        search_summary = []

        # Search in each discovered folder recursively
        for folder in discovered_folders:
            try:
                logger.info(f"Starting recursive search in folder: {folder}")

                # Recursively search this folder and all its subfolders
                folder_results = await search_folder_recursive(
                    oauth_client,
                    access_token,
                    site_id,
                    drive_id,
                    folder,
                    query,
                    config.org_name,
                    max_depth=2,  # Limit search depth to 2 levels
                    current_depth=0,
                    max_results=50,  # Limit total results to prevent excessive search time
                    max_folders_per_level=8,  # Limit folders per level
                )

                all_matching_documents.extend(folder_results["matching_docs"])
                search_summary.append(
                    {
                        "folder": folder,
                        "total_documents": folder_results["total_files"],
                        "matching_documents": len(folder_results["matching_docs"]),
                        "folders_searched": folder_results["folders_searched"],
                        "status": "success",
                    }
                )

                # Early termination if we have enough results across all folders
                if len(all_matching_documents) >= config.global_search_limit:
                    logger.info(
                        f"Global early termination: Found {len(all_matching_documents)} total results, stopping further folder searches"
                    )
                    break

            except Exception as e:
                logger.warning(f"Error searching folder {folder}: {e}")
                search_summary.append(
                    {
                        "folder": folder,
                        "total_documents": 0,
                        "matching_documents": 0,
                        "status": "error",
                        "error": str(e),
                        "message": f"Error accessing folder: {folder}",
                    }
                )

        return {
            "status": "success",
            "query": query,
            "matching_documents": all_matching_documents,
            "total_matches": len(all_matching_documents),
            "search_summary": search_summary,
            "folders_searched": discovered_folders,
            "organization": config.org_name,
            "message": f"Found {len(all_matching_documents)} matching documents in {config.org_name}",
        }

    except Exception as e:
        logger.error(f"Error in search_sharepoint_documents: {e}")
        return {
            "status": "error",
            "query": query,
            "error": str(e),
            "message": f"Error searching {config.org_name} SharePoint documents",
            "organization": config.org_name,
        }


@generic_sharepoint_server.tool()
async def list_sharepoint_folder_contents(
    folder_path: Optional[str] = None,
    include_subfolders: bool = True,
    user_token: Optional[str] = None,
) -> Dict[str, Any]:
    """
    List the contents of a SharePoint folder to explore structure using Microsoft Graph API.

    This tool helps explore SharePoint folder structures for any organization
    to understand document organization and get counts with user-scoped access.

    Args:
        folder_path: SharePoint folder path to explore (uses department default if not provided)
        include_subfolders: Whether to include subfolder information
        user_token: Optional user access token for delegated access

    Returns:
        JSON response with folder contents, subfolder counts, and file listings
    """
    if not config or not oauth_client:
        return {
            "status": "error",
            "error": "Server not initialized with department configuration",
            "message": "Please initialize server with a valid department prefix",
        }

    try:
        # Get user token from environment variable if not provided
        if not user_token:
            user_token = os.getenv("USER_JWT_TOKEN")

        # For local development, don't try OBO if no valid token available
        if user_token == "user_token_placeholder" or not user_token:
            logger.info(
                "No OAuth token available - using application-only authentication for local development"
            )
            user_token = None
        else:
            logger.info(
                f"Using OAuth token for SharePoint OBO flow: {bool(user_token)}"
            )

        # Use default folder if none specified
        if folder_path is None:
            folder_path = config.doc_library

        logger.info(
            f"Listing {config.org_name} SharePoint folder contents: {folder_path}"
        )

        # Get appropriate access token using OBO flow
        access_token = await oauth_client.get_access_token(user_token)
        if not access_token:
            return {
                "status": "error",
                "folder": folder_path,
                "error": "Failed to obtain access token",
                "message": f"Authentication failed for {config.org_name} SharePoint. Please ensure you have the necessary permissions and your session is valid.",
                "organization": config.org_name,
            }

        # Get site and drive information
        success, site_info = await oauth_client.get_site_info(access_token)
        if not success:
            return {
                "status": "error",
                "folder": folder_path,
                "error": site_info,
                "message": "Failed to access SharePoint site",
                "organization": config.org_name,
            }

        site_id = site_info.get("id")
        success, drives_data = await oauth_client.get_site_drives(access_token, site_id)
        if not success:
            return {
                "status": "error",
                "folder": folder_path,
                "error": drives_data,
                "message": "Failed to access document libraries",
                "organization": config.org_name,
            }

        drives = drives_data.get("value", [])
        if not drives:
            return {
                "status": "error",
                "folder": folder_path,
                "error": "No document libraries found",
                "message": "No accessible document libraries found",
                "organization": config.org_name,
            }

        drive_id = drives[0].get("id")

        # Get folder contents
        success, items_data = await oauth_client.get_drive_items(
            access_token, site_id, drive_id, folder_path
        )
        if not success:
            return {
                "status": "error",
                "folder": folder_path,
                "error": items_data,
                "message": "Failed to access folder contents",
                "organization": config.org_name,
            }

        items = items_data.get("value", [])

        # Organize contents
        folders = []
        files = []

        for item in items:
            if item.get("folder"):
                folder_info = {
                    "name": item.get("name"),
                    "id": item.get("id"),
                    "lastModified": item.get("lastModifiedDateTime"),
                    "webUrl": item.get("webUrl"),
                    "childCount": item.get("folder", {}).get("childCount", 0),
                    "type": "folder",
                    "organization": config.org_name,
                }
                folders.append(folder_info)
            else:
                file_info = {
                    "name": item.get("name"),
                    "id": item.get("id"),
                    "size": item.get("size"),
                    "lastModified": item.get("lastModifiedDateTime"),
                    "webUrl": item.get("webUrl"),
                    "fileType": (
                        item.get("name", "").split(".")[-1].lower()
                        if "." in item.get("name", "")
                        else "unknown"
                    ),
                    "type": "file",
                    "organization": config.org_name,
                }
                files.append(file_info)

        # Calculate summary statistics for projects/organizational items
        org_folders = [
            f
            for f in folders
            if any(
                keyword in f["name"].lower()
                for keyword in [
                    "project",
                    "prospect",
                    "initiative",
                    "program",
                    "strategy",
                ]
            )
        ]

        return {
            "status": "success",
            "folder": folder_path,
            "total_items": len(items),
            "folder_count": len(folders),
            "file_count": len(files),
            "folders": folders,
            "files": files,
            "organizational_folders": org_folders,
            "potential_project_count": sum(f.get("childCount", 0) for f in org_folders),
            "organization": config.org_name,
            "authentication": "Microsoft Graph API with On-Behalf-Of flow",
            "message": f"Found {len(folders)} folders and {len(files)} files in {config.org_name} {folder_path or 'root'}. Organizational folders: {len(org_folders)}",
        }

    except Exception as e:
        logger.error(f"Error listing folder contents: {e}")
        return {
            "status": "error",
            "folder": folder_path,
            "error": str(e),
            "message": f"Failed to list {config.org_name} folder contents",
            "organization": config.org_name,
        }


@generic_sharepoint_server.tool()
async def get_sharepoint_document_content(
    folder_name: str, file_name: str, user_token: Optional[str] = None
) -> Dict[str, Any]:
    """
    Retrieve the content of a specific SharePoint document using Microsoft Graph API.

    Args:
        folder_name: SharePoint folder containing the document
        file_name: Name of the document file
        user_token: Optional user access token for delegated access

    Returns:
        JSON response with document content and metadata
    """
    if not config or not oauth_client:
        return {
            "status": "error",
            "error": "Server not initialized with department configuration",
            "message": "Please initialize server with a valid department prefix",
        }

    try:
        # Get user token from environment variable if not provided
        if not user_token:
            user_token = os.getenv("USER_JWT_TOKEN")

        # For local development, don't try OBO if no valid token available
        if user_token == "user_token_placeholder" or not user_token:
            user_token = None

        # STRICT AUTHENTICATION CHECK: If delegated access is enabled, user token is REQUIRED
        if config.use_delegated_access:
            if not user_token or not user_token.strip():
                return {
                    "status": "error",
                    "error": "Authentication required",
                    "folder": folder_name,
                    "file": file_name,
                    "message": DELEGATED_ACCESS_AUTH_ERROR,
                    "organization": config.org_name,
                    "delegated_access_enabled": True,
                    "authentication_failed": True,
                }
        # No verbose logging during parallel processing

        # Silent document content retrieval
        # Get appropriate access token based on configuration
        if config.use_delegated_access:
            access_token = await oauth_client.get_access_token(user_token)
        else:
            access_token = await oauth_client.get_application_token()
        if not access_token:
            return {
                "status": "error",
                "folder": folder_name,
                "file": file_name,
                "error": "Failed to obtain access token",
                "message": f"Authentication failed for {config.org_name} SharePoint. Please ensure you have the necessary permissions to access this file.",
                "organization": config.org_name,
            }

        # Get site and drive information
        success, site_info = await oauth_client.get_site_info(access_token)
        if not success:
            return {
                "status": "error",
                "folder": folder_name,
                "file": file_name,
                "error": site_info,
                "message": "Failed to access SharePoint site",
                "organization": config.org_name,
            }

        site_id = site_info.get("id")
        success, drives_data = await oauth_client.get_site_drives(access_token, site_id)
        if not success:
            return {
                "status": "error",
                "folder": folder_name,
                "file": file_name,
                "error": drives_data,
                "message": "Failed to access document libraries",
                "organization": config.org_name,
            }

        drives = drives_data.get("value", [])
        if not drives:
            return {
                "status": "error",
                "folder": folder_name,
                "file": file_name,
                "error": "No document libraries found",
                "message": "No accessible document libraries found",
                "organization": config.org_name,
            }

        drive_id = drives[0].get("id")

        # Find the specific document
        success, items_data = await oauth_client.get_drive_items(
            access_token, site_id, drive_id, folder_name
        )
        if not success:
            return {
                "status": "error",
                "folder": folder_name,
                "file": file_name,
                "error": items_data,
                "message": f"Failed to access folder: {folder_name}",
                "organization": config.org_name,
            }

        items = items_data.get("value", [])
        target_file = None

        for item in items:
            if item.get("name") == file_name and item.get("file"):
                target_file = item
                break

        if not target_file:
            return {
                "status": "error",
                "folder": folder_name,
                "file": file_name,
                "error": "File not found",
                "message": f"Document '{file_name}' not found in folder '{folder_name}'",
                "organization": config.org_name,
            }

        # Download the file content
        file_id = target_file.get("id")
        success, file_content = await oauth_client.get_file_content(
            access_token, site_id, drive_id, file_id
        )

        if not success:
            return {
                "status": "error",
                "folder": folder_name,
                "file": file_name,
                "error": file_content,
                "message": "Failed to download file content",
                "organization": config.org_name,
            }

        # Parse the document content using embedded parser
        try:
            # Suppress any parsing logs
            import logging

            mpo_logger = logging.getLogger("MPO-SharePoint")
            original_level = mpo_logger.level
            mpo_logger.setLevel(logging.ERROR)  # Suppress extraction logs

            try:
                parsed_content = await parse_document_content_embedded(
                    file_content, file_name
                )
            finally:
                mpo_logger.setLevel(original_level)  # Restore log level

            return {
                "status": "success",
                "folder": folder_name,
                "file": file_name,
                "content": parsed_content,
                "size": target_file.get("size"),
                "lastModified": target_file.get("lastModifiedDateTime"),
                "webUrl": target_file.get("webUrl"),
                "organization": config.org_name,
                "authentication": "Microsoft Graph API with On-Behalf-Of flow",
                "message": f"Successfully extracted content from {file_name} in {config.org_name}",
            }

        except Exception as parse_error:
            logger.error(f"Error parsing document {file_name}: {parse_error}")
            return {
                "status": "error",
                "folder": folder_name,
                "file": file_name,
                "error": str(parse_error),
                "message": f"Failed to parse document content: {parse_error}",
                "organization": config.org_name,
            }

    except Exception as e:
        logger.error(f"Error getting document content: {e}")
        return {
            "status": "error",
            "folder": folder_name,
            "file": file_name,
            "error": str(e),
            "message": f"Failed to retrieve {config.org_name} document content",
            "organization": config.org_name,
        }


@generic_sharepoint_server.tool()
async def get_all_documents_comprehensive(
    user_token: Optional[str] = None,
) -> Dict[str, Any]:
    """
    Get ALL documents from SharePoint by traversing every single folder.
    This bypasses the broken search API and actually finds everything.
    """
    if not config or not oauth_client:
        return {
            "status": "error",
            "error": "Server not initialized",
            "message": "Server not initialized with department configuration",
        }

    try:
        # Silent comprehensive folder traversal
        # Get authentication token - check environment variable if not provided
        effective_token = user_token or os.getenv("USER_JWT_TOKEN")

        # STRICT AUTHENTICATION CHECK: If delegated access is enabled, user token is REQUIRED
        if config.use_delegated_access:
            if (
                not effective_token
                or effective_token == "user_token_placeholder"
                or not effective_token.strip()
            ):
                return {
                    "status": "error",
                    "error": "Authentication required",
                    "message": DELEGATED_ACCESS_AUTH_ERROR,
                    "organization": config.org_name,
                    "delegated_access_enabled": True,
                    "authentication_failed": True,
                }

        if config.use_delegated_access:
            access_token = await oauth_client.get_access_token(effective_token)
        else:
            access_token = await oauth_client.get_application_token()

        if not access_token:
            return {"status": "error", "message": "Failed to get access token"}

        # Get site and drives
        success, site_info = await oauth_client.get_site_info(access_token)
        if not success:
            return {
                "status": "error",
                "message": f"Failed to get site info: {site_info}",
            }

        site_id = site_info.get("id")
        success, drives_data = await oauth_client.get_site_drives(access_token, site_id)
        if not success:
            return {
                "status": "error",
                "message": f"Failed to get drives: {drives_data}",
            }

        drives = drives_data.get("value", [])
        if not drives:
            return {"status": "error", "message": "No drives found"}

        drive_id = drives[0].get("id")

        # Recursively get documents - limit to default_search_folders if configured
        all_documents = []
        folders_processed = {"count": 0}  # Track folder count

        # If default_search_folders is configured, only search those specific folders
        if config.default_search_folders:
            logger.info(f"Searching specific folders: {config.default_search_folders}")
            for folder_path in config.default_search_folders:
                await _traverse_all_folders_recursive(
                    oauth_client,
                    access_token,
                    site_id,
                    drive_id,
                    folder_path,
                    all_documents,
                    folders_processed,
                    max_depth=10,
                )
            search_scope = f"folders: {', '.join(config.default_search_folders)}"
        else:
            # No specific folders configured - search everything from root
            await _traverse_all_folders_recursive(
                oauth_client,
                access_token,
                site_id,
                drive_id,
                "",
                all_documents,
                folders_processed,
                max_depth=10,
            )
            search_scope = "all folders (no default_search_folders configured)"

        return {
            "status": "success",
            "total_documents": len(all_documents),
            "folders_processed": folders_processed["count"],
            "documents": all_documents,
            "organization": config.org_name,
            "search_scope": search_scope,
            "message": f"Found {len(all_documents)} documents from {folders_processed['count']} folders in {search_scope}",
        }

    except Exception as e:
        logger.error(f"Error in comprehensive document retrieval: {e}")
        return {
            "status": "error",
            "message": f"Failed to get all documents: {str(e)}",
            "organization": config.org_name,
        }


async def _traverse_all_folders_recursive(
    oauth_client,
    access_token: str,
    site_id: str,
    drive_id: str,
    folder_path: str,
    all_documents: list,
    folders_processed: dict,
    max_depth: int = 10,
    current_depth: int = 0,
):
    """Recursively traverse ALL folders and collect ALL documents"""
    if current_depth >= max_depth:
        return

    try:
        folders_processed["count"] += 1  # Increment folder count

        # Get all items in current folder
        success, items_data = await oauth_client.get_drive_items(
            access_token, site_id, drive_id, folder_path
        )

        if not success:
            logger.warning(f"Failed to access folder: {folder_path}")
            return

        items = items_data.get("value", [])

        for item in items:
            if item.get("file"):  # It's a file
                doc_info = {
                    "name": item.get("name", "Unknown"),
                    "id": item.get("id", ""),
                    "webUrl": item.get("webUrl", ""),
                    "lastModified": item.get("lastModifiedDateTime", ""),
                    "size": item.get("size", 0),
                    "folder_path": folder_path,
                    "organization": config.org_name,
                }
                all_documents.append(doc_info)

            elif item.get("folder"):  # It's a folder, recurse into it
                subfolder_name = item.get("name", "")
                if subfolder_name:
                    subfolder_path = (
                        f"{folder_path}/{subfolder_name}"
                        if folder_path
                        else subfolder_name
                    )
                    await _traverse_all_folders_recursive(
                        oauth_client,
                        access_token,
                        site_id,
                        drive_id,
                        subfolder_path,
                        all_documents,
                        folders_processed,
                        max_depth,
                        current_depth + 1,
                    )

    except Exception as e:
        logger.warning(f"Error traversing folder {folder_path}: {e}")


@generic_sharepoint_server.tool()
async def analyze_all_documents_for_content(
    search_terms: str,
    user_token: Optional[str] = None,
) -> Dict[str, Any]:
    """
    Get ALL documents and analyze each one for the search terms using highly optimized parallel processing.
    This uses concurrent processing with minimal logging for maximum speed.
    """
    if not config or not oauth_client:
        return {
            "status": "error",
            "error": "Server not initialized",
            "message": "Server not initialized with department configuration",
        }

    try:
        import time

        start_time = time.time()

        # Get user token from parameter or environment variable
        effective_token = user_token or os.getenv("USER_JWT_TOKEN")

        # STRICT AUTHENTICATION CHECK: If delegated access is enabled, user token is REQUIRED
        if config.use_delegated_access:
            if (
                not effective_token
                or effective_token == "user_token_placeholder"
                or not effective_token.strip()
            ):
                return {
                    "status": "error",
                    "error": "Authentication required",
                    "message": DELEGATED_ACCESS_AUTH_ERROR,
                    "organization": config.org_name,
                    "delegated_access_enabled": True,
                    "authentication_failed": True,
                }

        # Get all documents using comprehensive traversal (no verbose logging)
        all_docs_result = await get_all_documents_comprehensive(effective_token)

        if all_docs_result.get("status") != "success":
            return all_docs_result

        all_documents = all_docs_result.get("documents", [])
        total_docs = len(all_documents)
        folders_processed = all_docs_result.get("folders_processed", 0)

        # Use highly parallel processing for much faster analysis with crash prevention
        try:
            relevant_documents = await _analyze_documents_parallel_ultra_fast(
                all_documents, search_terms, effective_token
            )

            # Ensure we have a valid list
            if not isinstance(relevant_documents, list):
                relevant_documents = []

        except Exception as analysis_error:
            logger.error(f"Document analysis crashed: {str(analysis_error)}")
            return {
                "status": "error",
                "message": f"Document analysis failed: {str(analysis_error)}",
                "organization": config.org_name,
                "crash_prevention": "Applied",
            }

        # Calculate performance metrics
        end_time = time.time()
        processing_time = end_time - start_time

        # Sort by relevance score (highest first)
        relevant_documents.sort(key=lambda x: x["relevance_score"], reverse=True)

        # Simple response size limiting (remove double truncation that causes slowdown)
        limited_results = relevant_documents[:25]

        # Light content limiting - increased preview size to capture early pages of documents
        for doc in limited_results:
            if "content_preview" in doc:
                doc["content_preview"] = (
                    doc["content_preview"][:2000]
                    + "..."  # Increased to 2000 chars for multi-page content
                    if len(doc["content_preview"]) > 2000
                    else doc["content_preview"]
                )
            if "matches" in doc:
                doc["matches"] = doc["matches"][:7]  # Increased to 7 matches

        # FORCE performance logging - short message to prevent truncation
        perf_message = f"🚀 {folders_processed}F/{total_docs}D in {round(processing_time, 1)}s → {len(relevant_documents)} results"
        logger.info(perf_message)
        print(perf_message)

        # Also log to stderr for guaranteed visibility
        import sys

        print(perf_message, file=sys.stderr)

        # Build response with aggressive size control
        response = {
            "status": "success",
            "search_terms": search_terms,
            "documents_analyzed": total_docs,
            "relevant_documents": limited_results,
            "total_relevant": len(relevant_documents),
            "organization": config.org_name,
            "processing_time_seconds": round(processing_time, 1),
            "message": f"🚀 ANALYSIS COMPLETE: {folders_processed} folders, {total_docs} docs in {round(processing_time, 1)}s → {len(relevant_documents)} results (top {len(limited_results)})",
        }

        # Simple, fast response - no more double truncation

        # FINAL PERFORMANCE SUMMARY - appears after all processing
        final_summary = f"\n🚀 FINAL: {folders_processed} folders, {total_docs} documents → {round(processing_time, 1)}s total → {len(relevant_documents)} results"
        print(final_summary)
        logger.info(final_summary.strip())

        return response

    except Exception as e:
        return {
            "status": "error",
            "message": f"Analysis failed: {str(e)}",
            "organization": config.org_name,
        }


async def _analyze_documents_parallel_ultra_fast(
    documents: List[Dict], search_terms: str, user_token: Optional[str] = None
) -> List[Dict[str, Any]]:
    """Ultra-fast parallel document analysis with minimal logging and optimized performance."""
    relevant_documents = []
    max_workers = min(
        20, max(8, len(documents) // 3)
    )  # AGGRESSIVE: Up to 20 workers for maximum speed

    # Global content cache with locks
    content_cache = {}
    cache_hits = 0

    async def _process_document_fast(doc_info: Dict) -> Optional[Dict]:
        """Fast document processing with optimized content retrieval."""
        nonlocal cache_hits

        try:
            doc_name = doc_info.get("name", "Unknown")
            folder_path = doc_info.get("folder_path", "")

            # SPEED BOOST: Skip obviously irrelevant files by extension only
            filename_lower = doc_name.lower()
            search_words_lower = [w.lower() for w in search_terms.split() if len(w) > 2]

            # Skip only non-document files (keep all actual documents for content analysis)
            if not any(
                ext in filename_lower
                for ext in [".pdf", ".docx", ".doc", ".xlsx", ".pptx", ".txt"]
            ):
                return None  # Skip only non-document files

            # Quick filename relevance check
            search_lower = search_terms.lower()
            filename_score = 20 if search_lower in doc_name.lower() else 0

            # Check cache to avoid re-downloading same document within this query
            cache_key = f"{folder_path}::{doc_name}"
            if cache_key in content_cache:
                content = content_cache[cache_key]
                cache_hits += 1
            else:
                # Get document content - SILENT (no OAuth logging)
                import logging

                oauth_logger = logging.getLogger("sharepoint_oauth")
                mpo_logger = logging.getLogger("MPO-SharePoint")
                original_level_oauth = oauth_logger.level
                original_level_mpo = mpo_logger.level
                oauth_logger.setLevel(logging.ERROR)  # Suppress OAuth logs
                mpo_logger.setLevel(logging.ERROR)  # Suppress MPO logs

                try:
                    content_result = await get_sharepoint_document_content(
                        folder_path, doc_name, user_token
                    )
                    content = (
                        content_result.get("content", "")
                        if content_result.get("status") == "success"
                        else ""
                    )
                    content_cache[cache_key] = content
                finally:
                    # Restore original log levels
                    oauth_logger.setLevel(original_level_oauth)
                    mpo_logger.setLevel(original_level_mpo)

            if not content and filename_score == 0:
                return None

            # Fast content analysis
            content_lower = content.lower()
            search_words = [
                w.strip() for w in search_lower.split() if len(w.strip()) > 2
            ]

            # Score document relevance
            relevance_score = filename_score
            matches = []

            # Exact phrase match (highest priority)
            if search_lower in content_lower:
                relevance_score += 30
                # Get snippet around exact match
                idx = content_lower.find(search_lower)
                snippet = content[
                    max(0, idx - 40) : idx + len(search_terms) + 40
                ].strip()
                matches.append(snippet)

            # Individual word scoring
            for word in search_words:
                count = min(content_lower.count(word), 6)  # Cap count
                if count > 0:
                    relevance_score += count * 1.5

                    # Get one snippet per word (for speed)
                    if len(matches) < 4:
                        pos = content_lower.find(word)
                        if pos != -1:
                            snippet = content[
                                max(0, pos - 25) : pos + len(word) + 25
                            ].strip()
                            if snippet and snippet not in matches:
                                matches.append(snippet)

            # Return result if relevant
            if relevance_score > 0:
                return {
                    "document": {
                        "name": doc_info.get("name", "Unknown"),
                        "id": doc_info.get("id", ""),
                        "webUrl": doc_info.get("webUrl", ""),
                        "folder_path": doc_info.get("folder_path", ""),
                        "size": doc_info.get("size", 0),
                    },
                    "relevance_score": relevance_score,
                    "content_preview": _get_optimized_content_preview(content),
                    "matches": matches[:3],  # Reduced from 6 to 3
                    "organization": config.org_name,
                }

            return None

        except Exception:
            return None

    # Process documents with asyncio semaphore for concurrency control
    semaphore = asyncio.Semaphore(max_workers)

    async def _process_with_semaphore(doc_info):
        async with semaphore:
            try:
                # SPEED: Much shorter timeout for faster processing
                return await asyncio.wait_for(
                    _process_document_fast(doc_info),
                    timeout=15.0,  # Reduced to 15 seconds for speed
                )
            except asyncio.TimeoutError:
                return None  # Handle timeout gracefully
            except Exception:
                return None  # Silently handle any processing errors

    # Create all tasks and process them concurrently with robust error handling
    tasks = [_process_with_semaphore(doc_info) for doc_info in documents]

    # Process all tasks concurrently with progress tracking and error resilience
    results = []
    high_quality_count = 0

    try:
        # Use asyncio.as_completed for concurrent processing
        completed_count = 0
        for coro in asyncio.as_completed(tasks):
            try:
                result = await coro
                completed_count += 1

                if result:
                    relevant_documents.append(result)
                    if result["relevance_score"] > 15:
                        high_quality_count += 1

                    # AGGRESSIVE early termination for blazing speed
                    if high_quality_count >= 8 and len(relevant_documents) >= 12:
                        # Cancel remaining tasks IMMEDIATELY for speed
                        for task in tasks:
                            if not task.done():
                                try:
                                    task.cancel()
                                except Exception:
                                    pass
                        break  # Exit immediately

            except Exception:
                # Continue processing even if individual tasks fail
                continue

    except Exception:
        # If the entire processing fails, return what we have so far
        pass

    return relevant_documents


def _get_optimized_content_preview(content: str) -> str:
    """Ultra-fast content preview - no regex, just simple truncation for speed."""
    if not content:
        return ""

    # Much smaller preview to prevent MCP response size issues
    if len(content) <= 300:
        return content

    # Just take first 300 characters and add ellipsis
    return content[:300] + "..."


@generic_sharepoint_server.tool()
async def check_sharepoint_permissions(
    user_token: Optional[str] = None,
) -> Dict[str, Any]:
    """
    Check SharePoint connection status and user permissions.

    Args:
        user_token: Optional user access token for delegated access testing

    Returns:
        JSON response with connection status and permission details
    """
    if not config or not oauth_client:
        return {
            "status": "error",
            "error": "Server not initialized with department configuration",
            "message": "Please initialize server with a valid department prefix",
        }

    try:
        logger.info(
            f"Checking {config.org_name} SharePoint permissions and connection status"
        )

        # Get user token from parameter or environment variable
        effective_token = user_token or os.getenv("USER_JWT_TOKEN")

        # Test OAuth connection with OBO flow
        test_result = await oauth_client.test_connection(effective_token)

        return {
            "status": "success",
            "connection_test": test_result,
            "organization": config.org_name,
            "configuration": {
                "client_id": (
                    oauth_client.client_id[:8] + "..."
                    if oauth_client.client_id
                    else "missing"
                ),
                "tenant_id": oauth_client.tenant_id,
                "site_url": oauth_client.site_url,
                "delegated_access": oauth_client.use_delegated_access,
                "obo_scope": oauth_client.obo_scope,
                "default_folders": config.default_search_folders,
            },
            "authentication": "Microsoft Graph API with On-Behalf-Of flow",
            "message": f"{config.org_name} SharePoint permissions check completed",
        }

    except Exception as e:
        logger.error(f"Error checking {config.org_name} SharePoint permissions: {e}")
        return {
            "status": "error",
            "error": str(e),
            "message": f"Failed to check {config.org_name} SharePoint permissions",
            "organization": config.org_name,
        }


async def main(department_prefix: str = "MPO"):
    """Main function to run the SharePoint MCP server for a specific department"""
    try:
        # Initialize server for the specified department
        success = initialize_department_server(department_prefix)
        if not success:
            logger.error(
                f"Failed to initialize server for department: {department_prefix}"
            )
            return

        logger.info(f"Starting {config.org_name} SharePoint MCP Server")
        await mcp.run(transport="stdio")

    except Exception as e:
        logger.error(f"Error starting {department_prefix} SharePoint server: {e}")
        raise


if __name__ == "__main__":
    import sys

    # Allow department prefix to be passed as command line argument
    department_prefix = sys.argv[1] if len(sys.argv) > 1 else "MPO"

    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    )


# ===== EMBEDDED DOCUMENT PARSING FUNCTIONS =====


async def parse_document_content_embedded(content: bytes, file_name: str) -> str:
    """
    Parse document content using embedded parsers based on file type.

    Args:
        content: Raw document content bytes
        file_name: Name of the file to determine parsing method

    Returns:
        Extracted text content as string
    """
    try:
        file_ext = file_name.lower().split(".")[-1] if "." in file_name else ""

        # Parse silently - no logging during parallel processing
        if file_ext == "pdf":
            return await parse_pdf_content_embedded(content, file_name)
        elif file_ext in ["docx", "doc"]:
            return await parse_word_content_embedded(content, file_name)
        elif file_ext in ["pptx", "ppt"]:
            return await parse_pptx_content_embedded(content, file_name)
        elif file_ext in ["txt", "csv", "text"]:
            return await parse_text_content_embedded(content, file_name)
        else:
            return await parse_generic_content_embedded(content, file_name)

    except Exception as e:
        logger.error(f"Error parsing document {file_name}: {e}")
        return f"Error parsing document: {str(e)}"


async def parse_pdf_content_embedded(content: bytes, file_name: str) -> str:
    """
    Parse PDF content using pypdf and pdfplumber (embedded implementation).

    Args:
        content: PDF content bytes
        file_name: Name of the PDF file

    Returns:
        Extracted text content
    """
    try:
        from io import BytesIO

        # First try pypdf (if available)
        try:
            import pypdf

            pdf_stream = BytesIO(content)
            reader = pypdf.PdfReader(pdf_stream)

            text_content = []
            for page_num, page in enumerate(reader.pages):
                try:
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(f"--- Page {page_num + 1} ---")
                        text_content.append(text.strip())
                except Exception as page_error:
                    logger.warning(
                        f"Error extracting page {page_num + 1} from {file_name}: {page_error}"
                    )
                    continue

            if text_content:
                extracted_text = "\n\n".join(text_content)
                return extracted_text

        except (ImportError, Exception):
            pass  # Silent fallback

        # Fallback to pdfplumber
        try:
            import pdfplumber

            pdf_stream = BytesIO(content)
            text_content = []

            with pdfplumber.open(pdf_stream) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    try:
                        text = page.extract_text()
                        if text and text.strip():
                            text_content.append(f"--- Page {page_num + 1} ---")
                            text_content.append(text.strip())
                    except Exception:
                        continue  # Silent skip on page errors "\n\n".join(text_content)
                return extracted_text

        except Exception:
            pass  # Silent failure

        return f"Unable to extract text from PDF {file_name}. Both pypdf and pdfplumber failed."

    except Exception:
        return f"Error parsing PDF {file_name}"


async def parse_word_content_embedded(content: bytes, file_name: str) -> str:
    """
    Parse Word document content using python-docx (embedded implementation).

    Args:
        content: Word document content bytes
        file_name: Name of the Word file

    Returns:
        Extracted text content
    """
    try:
        from io import BytesIO
        import zipfile
        import xml.etree.ElementTree as ET

        # Parse DOCX as ZIP file and extract text from document.xml
        try:
            docx_stream = BytesIO(content)
            with zipfile.ZipFile(docx_stream, "r") as zip_file:
                # Read the main document content
                doc_xml = zip_file.read("word/document.xml")

                # Parse XML and extract text
                root = ET.fromstring(doc_xml)

                # Find all text elements (namespace handling)
                namespace = {
                    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                }
                text_elements = root.findall(".//w:t", namespace)

                text_content = []
                for elem in text_elements:
                    if elem.text:
                        text_content.append(elem.text)

                extracted_text = "".join(text_content)
                return extracted_text

        except Exception:
            return f"Error parsing Word document {file_name}"

    except Exception:
        return f"Error parsing Word document {file_name}"


async def parse_pptx_content_embedded(content: bytes, file_name: str) -> str:
    """
    Parse PowerPoint document content using python-pptx (embedded implementation).

    Args:
        content: PowerPoint document content bytes
        file_name: Name of the PowerPoint file

    Returns:
        Extracted text content from all slides
    """
    try:
        from io import BytesIO
        from pptx import Presentation

        # Parse PPTX
        pptx_stream = BytesIO(content)
        prs = Presentation(pptx_stream)

        text_content = []

        # Extract text from all slides
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            slide_text.append(f"--- Slide {slide_num} ---")

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if len(slide_text) > 1:  # Has more than just the slide number header
                text_content.append("\n".join(slide_text))

        if text_content:
            extracted_text = "\n\n".join(text_content)
            return extracted_text
        else:
            return f"No text content found in PowerPoint file {file_name}"

    except ImportError:
        return f"python-pptx library not available. Cannot parse {file_name}"
    except Exception as e:
        logger.warning(f"Error parsing PowerPoint document {file_name}: {e}")
        return f"Error parsing PowerPoint document {file_name}: {str(e)}"


async def parse_text_content_embedded(content: bytes, file_name: str) -> str:
    """
    Parse plain text content (embedded implementation).

    Args:
        content: Text content bytes
        file_name: Name of the text file

    Returns:
        Decoded text content
    """
    try:
        # Try different encodings
        encodings = ["utf-8", "utf-16", "latin-1", "cp1252"]

        for encoding in encodings:
            try:
                decoded_text = content.decode(encoding)
                return decoded_text.strip()
            except UnicodeDecodeError:
                continue

        # Silent fallback to utf-8 with error handling
        decoded_text = content.decode("utf-8", errors="ignore")
        return decoded_text

    except Exception:
        return f"Error parsing text file {file_name}"


async def parse_generic_content_embedded(content: bytes, file_name: str) -> str:
    """
    Parse generic/unknown content types (embedded implementation).

    Args:
        content: Content bytes
        file_name: Name of the file

    Returns:
        Best-effort text extraction
    """
    try:
        # Try to decode as text first
        try:
            decoded_text = content.decode("utf-8", errors="ignore")
            # Filter out non-printable characters for better readability
            import re

            cleaned_text = re.sub(
                r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\xff]", "", decoded_text
            )

            if len(cleaned_text.strip()) > 10:  # If we got reasonable text
                return cleaned_text
        except Exception as decode_error:
            logger.warning(f"Could not decode {file_name} as text: {decode_error}")

        # If text extraction failed, provide file info
        file_size = len(content)
        return f"Binary file {file_name} ({file_size} bytes). Content type not supported for text extraction."

    except Exception as e:
        logger.error(f"Error in generic content parsing for {file_name}: {e}")
        return f"Error parsing file {file_name}: {str(e)}"


if __name__ == "__main__":
    import sys

    # Allow department prefix to be passed as command line argument
    department_prefix = sys.argv[1] if len(sys.argv) > 1 else "MPO"

    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    )

    # Initialize server for the specified department
    success = initialize_department_server(department_prefix)
    if not success:
        logger.error(f"Failed to initialize server for department: {department_prefix}")
        sys.exit(1)

    logger.info(f"Starting {config.org_name} SharePoint MCP Server")
    # Run the server directly like other FastMCP servers
    mcp.run()
